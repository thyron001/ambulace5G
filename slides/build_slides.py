"""Presentación (4 partes) del proyecto ambulancia 5G — estilo minimalista."""
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TINTA = RGBColor(0x20, 0x20, 0x20)
GRIS = RGBColor(0x55, 0x55, 0x55)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def texto(slide, x, y, w, h, contenido, size, color=TINTA, bold=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(contenido.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def titulo(slide, t):
    texto(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.9), t, 26, TINTA, bold=True)
    ln = slide.shapes.add_connector(2, Inches(0.7), Inches(1.35), Inches(12.6), Inches(1.35))
    ln.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    ln.line.width = Pt(1)


def imagen(slide, ruta, top, max_w=Inches(11.8), max_h=Inches(5.5)):
    iw, ih = Image.open(ruta).size
    e = min(max_w / iw, max_h / ih)
    w, h = int(iw * e), int(ih * e)
    slide.shapes.add_picture(ruta, int((SW - w) / 2), top, width=w, height=h)


def s_seccion(num, t, subt=""):
    s = prs.slides.add_slide(BLANK)
    texto(s, Inches(1), Inches(2.9), Inches(11.3), Inches(0.7), f"PARTE {num}",
          18, GRIS, align=PP_ALIGN.CENTER)
    texto(s, Inches(1), Inches(3.5), Inches(11.3), Inches(1.5), t, 38, TINTA, bold=True,
          align=PP_ALIGN.CENTER)
    if subt:
        texto(s, Inches(1), Inches(4.9), Inches(11.3), Inches(0.8), subt, 18, GRIS,
              align=PP_ALIGN.CENTER)
    return s


def s_texto(t, cuerpo, size=20):
    s = prs.slides.add_slide(BLANK)
    titulo(s, t)
    texto(s, Inches(0.9), Inches(1.7), Inches(11.6), Inches(5.2), cuerpo, size, TINTA)
    return s


def s_imagen(t, ruta, pie=""):
    s = prs.slides.add_slide(BLANK)
    titulo(s, t)
    if pie:
        imagen(s, ruta, Inches(1.55), max_h=Inches(4.9))
        texto(s, Inches(0.9), Inches(6.7), Inches(11.6), Inches(0.7), pie, 13, GRIS,
              align=PP_ALIGN.CENTER)
    else:
        imagen(s, ruta, Inches(1.6), max_h=Inches(5.5))
    return s


# ---------- PORTADA ----------
s = prs.slides.add_slide(BLANK)
texto(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.6),
      "Ambulancia inteligente sobre red 5G", 38, TINTA, bold=True, align=PP_ALIGN.CENTER)
texto(s, Inches(1), Inches(4.1), Inches(11.3), Inches(1.2),
      "Co-simulación SUMO + ns-3 (5G-LENA)\nmMTC · URLLC + eMBB · Network slicing",
      20, GRIS, align=PP_ALIGN.CENTER)

# ---------- PARTE 1 ----------
s_seccion(1, "Arquitectura y parámetros de la red 5G")
s_imagen("Flujo del proyecto", "img/diagrama_flujo.jpg")

# tabla de parámetros
s = prs.slides.add_slide(BLANK)
titulo(s, "Parámetros generales de la red 5G")
params = [
    ("Parámetro", "Valor"),
    ("Frecuencia portadora", "3.5 GHz (banda n78, FR1)"),
    ("Ancho de banda", "10 – 20 MHz"),
    ("Numerología (SCS)", "1  (30 kHz)"),
    ("Duplexación", "TDD"),
    ("Modelo de propagación", "3GPP Urban Macro (UMa)"),
    ("Potencia Tx  (gNB / UE)", "43 dBm / 23 dBm"),
    ("Antenas  (gNB / UE)", "4×8 / 2×4  (MIMO)"),
    ("Scheduler", "OFDMA Round-Robin"),
    ("Slicing / QoS", "2 bandwidth parts:  URLLC / eMBB"),
    ("Handover", "X2  (multi-celda)"),
    ("Núcleo de red", "EPC (5G-LENA)"),
]
tbl = s.shapes.add_table(len(params), 2, Inches(1.6), Inches(1.7),
                         Inches(10.1), Inches(5.2)).table
tbl.columns[0].width = Inches(4.4)
tbl.columns[1].width = Inches(5.7)
for r in range(len(params)):
    for c in range(2):
        cell = tbl.cell(r, c)
        cell.text = params[r][c]
        pr = cell.text_frame.paragraphs[0].runs[0]
        pr.font.size = Pt(14)
        pr.font.bold = (r == 0)
        pr.font.color.rgb = BLANCO if r == 0 else TINTA
        cell.fill.solid()
        cell.fill.fore_color.rgb = (RGBColor(0x40, 0x40, 0x40) if r == 0
                                    else BLANCO)

# ---------- PARTE 2: mMTC ----------
s_seccion(2, "mMTC", "Massive Machine Type Communication")
s_texto("mMTC — comunicación masiva de máquinas",
    "Se consideran todos los vehículos de la ciudad en tiempo real para:\n\n"
    "•  Definir dinámicamente la ruta más corta de la ambulancia según el tráfico real.\n\n"
    "•  Poner los semáforos en verde a lo largo de la ruta (preempción), despejando el paso.\n\n"
    "•  El control se hace con TraCI sobre SUMO, con tráfico real y un accidente que "
    "obliga a re-rutear.")
s_imagen("Ruta generada por SUMO + TraCI (con accidente y tráfico real)", "img/mapa1.jpg",
         "La ambulancia recalcula su ruta al detectar el accidente; los gNBs son las celdas que la sirven.")
s_imagen("Conectividad 5G a lo largo de la ruta (SINR vs. tiempo)", "img/sinr.jpg",
         "Con varias celdas y handover, el SINR no cae a cero: la conexión se mantiene todo el trayecto.")
s_imagen("Efecto de la asistencia 5G en el tiempo de viaje", "img/sistema.jpg",
         "El re-ruteo dinámico y la preempción de semáforos reducen el tiempo de viaje y las detenciones.")

# ---------- PARTE 3 ----------
s_seccion(3, "Generación de tráfico 5G")
s_imagen("Filtrado de los gNBs que cruza la ruta", "img/cobertura.jpg",
         "Toda la ciudad tiene cobertura 5G (35 gNBs); se filtran las celdas por las que pasa la ambulancia.")
s_imagen("Autos que compartieron celda con la ambulancia", "img/competidores.jpg",
         "Se filtran los autos que estuvieron en esas celdas en los mismos instantes; sus rutas generan el tráfico en ns-3.")

# ---------- PARTE 4 ----------
s_seccion(4, "URLLC + eMBB")
s_texto("La ambulancia transmite al hospital por 5G",
    "La ambulancia envía datos al servidor del hospital por la red 5G:\n\n"
    "•  Video en tiempo real  →  eMBB\n"
    "•  Signos vitales  →  URLLC\n\n"
    "Se hicieron 3 simulaciones en ns-3:\n"
    "   1)  Solo tráfico de la ambulancia.\n"
    "   2)  Ambulancia + autos, sin slicing.\n"
    "   3)  Ambulancia + autos, con slicing.\n\n"
    "El video de la ambulancia y el tráfico de los autos son ambos eMBB.")
s_texto("¿Qué es el network slicing?",
    "El slicing divide los recursos de radio en rebanadas independientes, una por servicio:\n\n"
    "•  El espectro se parte en 2 bandas dedicadas:\n"
    "      –  una para URLLC (signos vitales), aislada.\n"
    "      –  otra para eMBB (video + tráfico de los autos).\n\n"
    "•  Aunque la banda eMBB se sature, el tráfico crítico URLLC viaja en recursos "
    "separados y queda protegido.")
s_imagen("Resultado: URLLC vs eMBB en las 3 simulaciones", "img/fase3.jpg",
         "Sin slicing se pierde el tráfico crítico bajo carga; con slicing los signos vitales se entregan intactos.")
s_imagen("Scheduler — throughput por flujo en el tiempo", "img/sched_thr.jpg",
         "El Round-Robin reparte la celda entre ambulancia y autos; con slicing los vitales ocupan una franja aislada.")
s_imagen("Scheduler — reparto de recursos por flujo", "img/sched_rep.jpg",
         "Sin slicing el reparto es parejo; con slicing los signos vitales tienen su porción reservada.")

out = "/home/thyron001/Desktop/moviles/slides/presentacion_ambulancia_5G.pptx"
prs.save(out)
print(f"Guardado {out}  ({len(prs.slides._sldIdLst)} slides)")
